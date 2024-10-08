import numpy as np
import pandas as pd

import json
import matplotlib
import matplotlib.pyplot as plt
import matplotlib.ticker as tck
from cycler import cycler
from mpl_toolkits.mplot3d import Axes3D
import seaborn as sns

import sys, os, shutil, re, yaml, json, subprocess
import importlib.util
import time
from dateutil import parser
from datetime import datetime

from PIL import Image, ImageEnhance, ImageOps, ImageFilter
from rembg import remove, new_session

import docx
from fpdf import FPDF
from lxml import etree
from docx import Document
from PyPDF2 import PdfReader
from pptx import Presentation
from pptx.util import Inches
from pdf2image import convert_from_path, pdfinfo_from_path
from nltk.tokenize import sent_tokenize, word_tokenize
import nltk  # nltk.download("punkt")
from docx2pdf import convert
import img2pdf as image2pdf
import nbformat
from nbconvert import MarkdownExporter

from itertools import pairwise
from box import Box, BoxList
from numerizer import numerize
from tqdm import tqdm
import mimetypes
from pprint import pp
from collections import Counter
from fuzzywuzzy import fuzz, process
from langdetect import detect
from duckduckgo_search import DDGS

from . import netfinder

# from .plot import get_color

try:
    get_ipython().run_line_magic("load_ext", "autoreload")
    get_ipython().run_line_magic("autoreload", "2")
except NameError:
    pass


def is_package_installed(package_name):
    """Check if a package is installed."""
    package_spec = importlib.util.find_spec(package_name)
    return package_spec is not None


def upgrade(module="py2ls"):
    # try:
    #     pkg_resources.get_distribution(module)
    # except pkg_resources.DistributionNotFound:
    #     subprocess.check_call([sys.executable, "-m", "pip", "install", module])
    if not is_package_installed(module):
        try:
            subprocess.check_call([sys.executable, "-m", "pip", "install", module])
        except subprocess.CalledProcessError as e:
            print(f"An error occurred while installing {module}: {e}")
    try:
        subprocess.check_call(
            [sys.executable, "-m", "pip", "install", "--upgrade", module]
        )
    except subprocess.CalledProcessError as e:
        print(f"An error occurred while upgrading py2ls: {e}")


dir_save = "/Users/macjianfeng/Dropbox/Downloads/"


def get_version(pkg):
    import importlib.metadata

    def get_v(pkg_name):
        try:
            version = importlib.metadata.version(pkg_name)
            print(f"version {pkg_name} == {version}")
        except importlib.metadata.PackageNotFoundError:
            print(f"Package '{pkg_name}' not found")

    if isinstance(pkg, str):
        get_v(pkg)
    elif isinstance(pkg, list):
        [get_v(pkg_) for pkg_ in pkg]


# usage:
# get_version(['pandas','numpy','py2ls'])


def rm_folder(folder_path, verbose=True):
    try:
        shutil.rmtree(folder_path)
        if verbose:
            print(f"Successfully deleted {folder_path}")
    except Exception as e:
        if verbose:
            print(f"Failed to delete {folder_path}. Reason: {e}")


def fremove(path, verbose=True):
    """
    Remove a folder and all its contents or a single file.
    Parameters:
    path (str): The path to the folder or file to remove.
    verbose (bool): If True, print success or failure messages. Default is True.
    """
    try:
        if os.path.isdir(path):
            shutil.rmtree(path)
            if verbose:
                print(f"Successfully deleted folder {path}")
        elif os.path.isfile(path):
            os.remove(path)
            if verbose:
                print(f"Successfully deleted file {path}")
        else:
            if verbose:
                print(f"Path {path} does not exist")
    except Exception as e:
        if verbose:
            print(f"Failed to delete {path}. Reason: {e}")


def get_cwd(verbose: bool = True):
    """
    get_cwd: to get the current working directory
    Args:
        verbose (bool, optional): to show which function is use. Defaults to True.
    """
    try:
        script_dir = os.path.dirname(os.path.abspath(__file__))
        if verbose:
            print("os.path.dirname(os.path.abspath(__file__)):", script_dir)
    except NameError:
        # This works in an interactive environment (like a Jupyter notebook)
        script_dir = os.getcwd()
        if verbose:
            print("os.getcwd():", script_dir)
    return script_dir


def search(
    query,
    limit=5,
    kind="text",
    output="df",
    verbose=False,
    download=False,
    dir_save=dir_save,
    **kwargs,
):

    if "te" in kind.lower():
        results = DDGS().text(query, max_results=limit)
        res = pd.DataFrame(results)
        res.rename(columns={"href": "links"}, inplace=True)
    if verbose:
        print(f'searching "{query}": got the results below\n{res}')
    if download:
        try:
            downloader(
                url=res.links.tolist(), dir_save=dir_save, verbose=verbose, **kwargs
            )
        except:
            if verbose:
                print(f"failed link")
    return res


def echo(*args, **kwargs):
    """
    query, model="gpt", verbose=True, log=True, dir_save=dir_save
    a ai chat tool
    Args:
        query (str): _description_
        model (str, optional): _description_. Defaults to "gpt".
        verbose (bool, optional): _description_. Defaults to True.
        log (bool, optional): _description_. Defaults to True.
        dir_save (str, path, optional): _description_. Defaults to dir_save.

    Returns:
        str: the answer from ai
    """
    global dir_save

    query = None
    model = kwargs.get("model", "gpt")
    verbose = kwargs.get("verbose", True)
    log = kwargs.get("log", True)
    dir_save = kwargs.get("dir_save", dir_save)
    for arg in args:
        if isinstance(arg, str):
            if os.path.isdir(arg):
                dir_save = arg
            # elif os.path.isfile(arg):
            #     dir_save = dirname(arg)
            elif len(arg) <= 5:
                model = arg
            else:
                query = arg
        elif isinstance(arg, dict):
            verbose = arg.get("verbose", verbose)
            log = arg.get("log", log)

    def is_in_any(str_candi_short, str_full, ignore_case=True):
        if isinstance(str_candi_short, str):
            str_candi_short = [str_candi_short]
        res_bool = []
        if ignore_case:
            [res_bool.append(i in str_full.lower()) for i in str_candi_short]
        else:
            [res_bool.append(i in str_full) for i in str_candi_short]
        return any(res_bool)

    def valid_mod_name(str_fly):
        if is_in_any(str_fly, "claude-3-haiku"):
            return "claude-3-haiku"
        elif is_in_any(str_fly, "gpt-3.5"):
            return "gpt-3.5"
        elif is_in_any(str_fly, "llama-3-70b"):
            return "llama-3-70b"
        elif is_in_any(str_fly, "mixtral-8x7b"):
            return "mixtral-8x7b"
        else:
            print(
                f"not support your model{model}, supported models: 'claude','gpt(default)', 'llama','mixtral'"
            )
            return "gpt-3.5"  # default model

    model_valid = valid_mod_name(model)
    res = DDGS().chat(query, model=model_valid)
    if verbose:
        pp(res)
    if log:
        dt_str = datetime.fromtimestamp(time.time()).strftime("%Y-%m-%d_%H:%M:%S")
        res_ = f"\n\n####Q:{query}\n\n#####Ans:{dt_str}\n\n>{res}\n"
        if bool(os.path.basename(dir_save)):
            fpath = dir_save
        else:
            os.makedirs(dir_save, exist_ok=True)
            fpath = os.path.join(dir_save, f"log_ai.md")
        fupdate(fpath=fpath, content=res_)
        print(f"log file:{fpath}")
    return res


def chat(*args, **kwargs):
    return echo(*args, **kwargs)


def ai(*args, **kwargs):
    return echo(*args, **kwargs)


def detect_lang(text, output="lang", verbose=True):
    dir_curr_script = os.path.dirname(os.path.abspath(__file__))
    dir_lang_code = dir_curr_script + "/data/lang_code_iso639.json"
    print(dir_curr_script, os.getcwd(), dir_lang_code)
    lang_code_iso639 = fload(dir_lang_code)
    l_lang, l_code = [], []
    [[l_lang.append(v), l_code.append(k)] for v, k in lang_code_iso639.items()]
    try:
        if is_text(text):
            code_detect = detect(text)
            if "c" in output.lower():  # return code
                return l_code[strcmp(code_detect, l_code, verbose=verbose)[1]]
            else:
                return l_lang[strcmp(code_detect, l_code, verbose=verbose)[1]]
        else:
            print(f"{text} is not supported")
            return "no"
    except:
        return "no"


def is_text(s):
    has_alpha = any(char.isalpha() for char in s)
    has_non_alpha = any(not char.isalpha() for char in s)
    # no_special = not re.search(r'[^A-Za-z0-9\s]', s)
    return has_alpha and has_non_alpha


def strcmp(search_term, candidates, ignore_case=True, verbose=True, scorer="WR"):
    """
    Compares a search term with a list of candidate strings and finds the best match based on similarity score.

    Parameters:
    search_term (str): The term to be searched for.
    candidates (list of str): A list of candidate strings to compare against the search term.
    ignore_case (bool): If True, the comparison ignores case differences.
    verbose (bool): If True, prints the similarity score and the best match.

    Returns:
    tuple: A tuple containing the best match and its index in the candidates list.
    """

    def to_lower(s, ignore_case=True):
        # Converts a string or list of strings to lowercase if ignore_case is True.
        if ignore_case:
            if isinstance(s, str):
                return s.lower()
            elif isinstance(s, list):
                return [elem.lower() for elem in s]
        return s

    str1_, str2_ = to_lower(search_term, ignore_case), to_lower(candidates, ignore_case)
    if isinstance(str2_, list):
        if "part" in scorer.lower():
            similarity_scores = [fuzz.partial_ratio(str1_, word) for word in str2_]
        elif "W" in scorer.lower():
            similarity_scores = [fuzz.WRatio(str1_, word) for word in str2_]
        elif "Ratio" in scorer.lower():
            similarity_scores = [fuzz.Ratio(str1_, word) for word in str2_]
        else:
            similarity_scores = [fuzz.WRatio(str1_, word) for word in str2_]
        best_match_index = similarity_scores.index(max(similarity_scores))
        best_match_score = similarity_scores[best_match_index]
    else:
        best_match_index = 0
        if "part" in scorer.lower():
            best_match_score = fuzz.partial_ratio(str1_, str2_)
        elif "W" in scorer.lower():
            best_match_score = fuzz.WRatio(str1_, str2_)
        elif "Ratio" in scorer.lower():
            best_match_score = fuzz.Ratio(str1_, str2_)
        else:
            best_match_score = fuzz.WRatio(str1_, str2_)
    if verbose:
        print(f"\nbest_match is: {candidates[best_match_index],best_match_score}")
        best_match = process.extract(search_term, candidates)
        print(f"建议: {best_match}")
    return candidates[best_match_index], best_match_index


# Example usaged
# str1 = "plos biology"
# str2 = ['PLoS Computational Biology', 'PLOS BIOLOGY']
# best_match, idx = strcmp(str1, str2, ignore_case=1)


def counter(list_, verbose=True):
    c = Counter(list_)
    # Print the name counts
    for item, count in c.items():
        if verbose:
            print(f"{item}: {count}")
    return c


# usage:
# print(f"Return an iterator over elements repeating each as many times as its count:\n{sorted(c.elements())}")
# print(f"Return a list of the n most common elements:\n{c.most_common()}")
# print(f"Compute the sum of the counts:\n{c.total()}")


def str2time(time_str, fmt="24"):
    """
    Convert a time string into the specified format.
    Parameters:
    - time_str (str): The time string to be converted.
    - fmt (str): The format to convert the time to. Defaults to '%H:%M:%S'.
    Returns:
        %I represents the hour in 12-hour format.
        %H represents the hour in 24-hour format (00 through 23).
        %M represents the minute.
        %S represents the second.
        %p represents AM or PM.
    - str: The converted time string.
    """

    def time_len_corr(time_str):
        time_str_ = (
            ssplit(time_str, by=[":", " ", "digital_num"]) if ":" in time_str else None
        )
        time_str_split = []
        [time_str_split.append(i) for i in time_str_ if is_num(i)]
        if time_str_split:
            if len(time_str_split) == 2:
                H, M = time_str_split
                time_str_full = H + ":" + M + ":00"
            elif len(time_str_split) == 3:
                H, M, S = time_str_split
                time_str_full = H + ":" + M + ":" + S
        else:
            time_str_full = time_str_
        if "am" in time_str.lower():
            time_str_full += " AM"
        elif "pm" in time_str.lower():
            time_str_full += " PM"
        return time_str_full

    if "12" in fmt:
        fmt = "%I:%M:%S %p"
    elif "24" in fmt:
        fmt = "%H:%M:%S"

    try:
        # Try to parse the time string assuming it could be in 24-hour or 12-hour format
        time_obj = datetime.strptime(time_len_corr(time_str), "%H:%M:%S")
    except ValueError:
        try:
            time_obj = datetime.strptime(time_len_corr(time_str), "%I:%M:%S %p")
        except ValueError as e:
            raise ValueError(f"Unable to parse time string: {time_str}. Error: {e}")

    # Format the time object to the desired output format
    formatted_time = time_obj.strftime(fmt)
    return formatted_time


# # Example usage:
# time_str1 = "14:30:45"
# time_str2 = "02:30:45 PM"

# formatted_time1 = str2time(time_str1, fmt='12')  # Convert to 12-hour format
# formatted_time2 = str2time(time_str2, fmt='24')    # Convert to 24-hour format

# print(formatted_time1)  # Output: 02:30:45 PM
# print(formatted_time2)  # Output: 14:30:45


def str2date(date_str, fmt="%Y-%m-%d_%H:%M:%S"):
    """
    Convert a date string into the specified format.
    Parameters:
    - date_str (str): The date string to be converted.
    - fmt (str): The format to convert the date to. Defaults to '%Y%m%d'.
    Returns:
    - str: The converted date string.
    """
    try:
        date_obj = parser.parse(date_str)
    except ValueError as e:
        raise ValueError(f"Unable to parse date string: {date_str}. Error: {e}")
    # Format the date object to the desired output format
    formatted_date = date_obj.strftime(fmt)
    return formatted_date


# str1=str2date(num2str(20240625),fmt="%a %d-%B-%Y")
# print(str1)
# str2=str2num(str2date(str1,fmt='%a %Y%m%d'))
# print(str2)


def str2num(s, *args, **kwargs):
    delimiter = kwargs.get("sep", None)
    round_digits = kwargs.get("round", None)
    if delimiter is not None:
        s = s.replace(delimiter, "")
    for arg in args:
        if isinstance(arg, str) and delimiter is None:
            delimiter = arg
        elif isinstance(arg, int) and round_digits is None:
            round_digits = arg
    try:
        num = int(s)
    except ValueError:
        try:
            num = float(s)
        except ValueError:
            try:
                numerized = numerize(s)
                num = int(numerized) if "." not in numerized else float(numerized)
            except Exception as e:
                # Attempt to handle multiple number segments
                try:
                    number_segments = ssplit(s, by="number_strings")
                    nums = []
                    for segment in number_segments:
                        try:
                            nums.append(str2num(segment))
                        except ValueError:
                            continue
                    if len(nums) == 1:
                        num = nums[0]
                    else:
                        raise ValueError(
                            "Multiple number segments found, cannot determine single numeric value"
                        )
                except Exception as e:
                    return None

    # Apply rounding if specified
    if round_digits is not None:
        num_adj = num + 0.00000000001  # Ensure precise rounding
        num = round(num_adj, round_digits)
    if round_digits == 0:
        num = int(num)
    # if delimiter is not None:
    #     num_str = f"{num:,}".replace(",", delimiter)
    #     return num_str#s.replace(delimiter, "")

    return num


# Examples
# print(str2num("123"))                # Output: 123
# print(str2num("123.456", 2))         # Output: 123.46
# print(str2num("one hundred and twenty three"))  # Output: 123
# print(str2num("seven million"))      # Output: 7000000
# print(str2num('one thousand thirty one',','))  # Output: 1,031
# print(str2num("12345.6789", ","))    # Output: 12,345.6789
# print(str2num("12345.6789", " ", 2)) # Output: 12 345.68
# print(str2num('111113.34555',3,',')) # Output: 111,113.346
# print(str2num("123.55555 sec miniuets",3)) # Output: 1.3
def num2str(num, *args, **kwargs):
    delimiter = kwargs.get("sep", None)
    round_digits = kwargs.get("round", None)

    # Parse additional arguments
    for arg in args:
        if isinstance(arg, str):
            delimiter = arg
        elif isinstance(arg, int):
            round_digits = arg

    # Apply rounding if specified
    if round_digits is not None:
        num = round(num, round_digits)

    # Convert number to string
    num_str = f"{num}"

    # Apply delimiter if specified
    if delimiter is not None:
        num_str = num_str.replace(".", ",")  # Replace decimal point with comma
        num_str_parts = num_str.split(",")
        if len(num_str_parts) > 1:
            integer_part = num_str_parts[0]
            decimal_part = num_str_parts[1]
            integer_part = "{:,}".format(int(integer_part))
            num_str = integer_part + "." + decimal_part
        else:
            num_str = "{:,}".format(int(num_str_parts[0]))

    return num_str


# Examples
# print(num2str(123),type(num2str(123)))                # Output: "123"
# print(num2str(123.456, 2),type(num2str(123.456, 2)))         # Output: "123.46"
# print(num2str(7000.125, 2),type(num2str(7000.125, 2)))        # Output: "7000.13"
# print(num2str(12345.6789, ","),type(num2str(12345.6789, ",")))    # Output: "12,345.6789"
# print(num2str(7000.00, ","),type(num2str(7000.00, ",")))       # Output: "7,000.00"
def sreplace(*args, **kwargs):
    """
    sreplace(text, by=None, robust=True)
    Replace specified substrings in the input text with provided replacements.
    Args:
        text (str): The input text where replacements will be made.
        by (dict, optional): A dictionary containing substrings to be replaced as keys
            and their corresponding replacements as values. Defaults to {".com": "..come", "\n": " ", "\t": " ", "  ": " "}.
        robust (bool, optional): If True, additional default replacements for newline and tab characters will be applied.
                                Default is False.
    Returns:
        str: The text after replacements have been made.
    """
    text = None
    by = kwargs.get("by", None)
    robust = kwargs.get("robust", True)

    for arg in args:
        if isinstance(arg, str):
            text = arg
        elif isinstance(arg, dict):
            by = arg
        elif isinstance(arg, bool):
            robust = arg
        else:
            Error(f"{type(arg)} is not supported")

    # Default replacements for newline and tab characters
    default_replacements = {
        "\a": "",
        "\b": "",
        "\f": "",
        "\n": "",
        "\r": "",
        "\t": "",
        "\v": "",
        "\\": "",  # Corrected here
        # "\?": "",
        "�": "",
        "\\x": "",  # Corrected here
        "\\x hhhh": "",
        "\\ ooo": "",  # Corrected here
        "\xa0": "",
        "  ": " ",
    }

    # If dict_replace is None, use the default dictionary
    if by is None:
        by = {}
    # If robust is True, update the dictionary with default replacements
    if robust:
        by.update(default_replacements)

    # Iterate over each key-value pair in the dictionary and replace substrings accordingly
    for k, v in by.items():
        text = text.replace(k, v)
    return text


# usage:
# sreplace(text, by=dict(old_str='new_str'), robust=True)


def paper_size(paper_type_str="a4"):
    df = pd.DataFrame(
        {
            "a0": [841, 1189],
            "a1": [594, 841],
            "a2": [420, 594],
            "a3": [297, 420],
            "a4": [210, 297],
            "a5": [148, 210],
            "a6": [105, 148],
            "a7": [74, 105],
            "b0": [1028, 1456],
            "b1": [707, 1000],
            "b2": [514, 728],
            "b3": [364, 514],
            "b4": [257, 364],
            "b5": [182, 257],
            "b6": [128, 182],
            "letter": [215.9, 279.4],
            "legal": [215.9, 355.6],
            "business card": [85.6, 53.98],
            "photo china passport": [33, 48],
            "passport single": [125, 88],
            "visa": [105, 74],
            "sim": [25, 15],
        }
    )
    for name in df.columns:
        if paper_type_str in name.lower():
            paper_type = name
    if not paper_type:
        paper_type = "a4"  # default
    return df[paper_type].tolist()


def docx2pdf(dir_docx, dir_pdf=None):
    if dir_pdf:
        convert(dir_docx, dir_pdf)
    else:
        convert(dir_docx)


def img2pdf(dir_img, kind="jpeg", page=None, dir_save=None, page_size="a4", dpi=300):
    def mm_to_point(size):
        return (image2pdf.mm_to_pt(size[0]), image2pdf.mm_to_pt(size[1]))

    def set_dpi(x):
        dpix = dpiy = x
        return image2pdf.get_fixed_dpi_layout_fun((dpix, dpiy))

    if not kind.startswith("."):
        kind = "." + kind
    if dir_save is None:
        dir_save = dir_img.replace(kind, ".pdf")
    imgs = []
    if os.path.isdir(dir_img):
        if not dir_save.endswith(".pdf"):
            dir_save += "#merged_img2pdf.pdf"
        if page is None:
            select_range = listdir(dir_img, kind=kind).fpath
        else:
            if not isinstance(page, (np.ndarray, list, range)):
                page = [page]
            select_range = listdir(dir_img, kind=kind)["fpath"][page]
        for fname in select_range:
            if not fname.endswith(kind):
                continue
            path = os.path.join(dir_img, fname)
            if os.path.isdir(path):
                continue
            imgs.append(path)
    else:
        imgs = [os.path.isdir(dir_img), dir_img]

    if page_size:
        if isinstance(page_size, str):
            pdf_in_mm = mm_to_point(paper_size(page_size))
        else:
            print("default: page_size = (210,297)")
            pdf_in_mm = mm_to_point(page_size)
            print(f"page size was set to {page_size}")
        p_size = image2pdf.get_layout_fun(pdf_in_mm)
    else:
        p_size = set_dpi(dpi)
    with open(dir_save, "wb") as f:
        f.write(image2pdf.convert(imgs, layout_fun=p_size))


# usage:
# dir_img="/Users/macjianfeng/Dropbox/00-Personal/2015-History/2012-2015_兰州大学/120901-大学课件/生物统计学 陆卫/复习题/"
# img2pdf(dir_img,kind='tif', page=range(3,7,2))


def pdf2ppt(dir_pdf, dir_ppt):
    prs = Presentation()

    # Open the PDF file
    with open(dir_pdf, "rb") as f:
        reader = PdfReader(f)
        num_pages = len(reader.pages)

        # Iterate through each page in the PDF
        for page_num in range(num_pages):
            page = reader.pages[page_num]
            text = page.extract_text()

            # Add a slide for each page's content
            slide_layout = prs.slide_layouts[
                5
            ]  # Use slide layout that suits your needs
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = f"Page {page_num + 1}"
            slide.shapes.add_textbox(
                Inches(1), Inches(1.5), Inches(8), Inches(5)
            ).text = text

    # Save the PowerPoint presentation
    prs.save(dir_ppt)
    print(f"Conversion from {dir_pdf} to {dir_ppt} complete.")


def ssplit(text, by="space", verbose=False, strict=False, **kws):
    if isinstance(text, list):
        nested_list = [ssplit(i, by=by, verbose=verbose, **kws) for i in text]
        flat_list = [item for sublist in nested_list for item in sublist]
        return flat_list

    def split_by_word_length(text, length):
        return [word for word in text.split() if len(word) == length]

    def split_by_multiple_delimiters(text, delimiters):
        regex_pattern = "|".join(map(re.escape, delimiters))
        return re.split(regex_pattern, text)

    def split_by_camel_case(text):
        return re.findall(r"[A-Z](?:[a-z]+|[A-Z]*(?=[A-Z]|$))", text)

    def split_at_upper_fl_lower(text):
        return re.findall(r"[A-Z](?:[a-z]+|[A-Z]+(?=[A-Z]|$))", text)

    def split_at_lower_fl_upper(text):
        split_text = re.split(r"(?<=[a-z])(?=[A-Z])", text)
        return split_text

    def split_at_upper(text):
        split_text = re.split(r"(?=[A-Z])", text)
        split_text = [part for part in split_text if part]
        return split_text

    def split_by_regex_lookahead(text, pattern):
        return re.split(f"(?<={pattern})", text)

    def split_by_regex_end(text, pattern):
        return re.split(f"(?={pattern})", text)

    # def split_by_sentence_endings(text):
    #     return re.split(r"(?<=[.!?])", text)
    def split_non_ascii(text):
        # return re.split(r"([^\x00-\x7F\w\s,.!?:\"'()\-]+)", text)
        # return re.split(r"[^\x00-\x7F]+", text)
        return re.split(r"([^\x00-\x7F]+)", text)

    def split_by_consecutive_non_alphanumeric(text):
        return re.split(r"\W+", text)

    def split_by_fixed_length_chunks(text, length):
        return [text[i : i + length] for i in range(0, len(text), length)]

    def split_by_sent_num(text, n=10):
        # split text into sentences
        text_split_by_sent = sent_tokenize(text)
        cut_loc_array = np.arange(0, len(text_split_by_sent), n)
        if cut_loc_array[-1] != len(text_split_by_sent):
            cut_loc = np.append(cut_loc_array, len(text_split_by_sent))
        else:
            cut_loc = cut_loc_array
        # get text in section (e.g., every 10 sentences)
        text_section = []
        for i, j in pairwise(cut_loc):
            text_section.append(text_split_by_sent[i:j])
        return text_section

    def split_general(text, by, verbose=False, ignore_case=False):
        if ignore_case:
            if verbose:
                print(f"used {by} to split, ignore_case=True")
            pattern = re.compile(re.escape(by), re.IGNORECASE)
            split_text = pattern.split(text)
            return split_text
        else:
            if verbose:
                print(f"used {by} to split, ignore_case=False")
            return text.split(by)

    def reg_split(text, pattern):
        return re.split(pattern, text)

    if ("sp" in by or "white" in by) and not strict:
        if verbose:
            print(f"splited by space")
        return text.split()
    elif ("word" in by and "len" in by) and not strict:
        if verbose:
            print(f"split_by_word_length(text, length)")
        return split_by_word_length(text, **kws)  # split_by_word_length(text, length)
    # elif "," in by:
    #     if verbose:
    #         print(f"splited by ','")
    #     return text.split(",")
    elif isinstance(by, list):
        if verbose:
            print(f"split_by_multiple_delimiters: ['|','&']")
        return split_by_multiple_delimiters(text, by)
    elif (
        all([("digi" in by or "num" in by), not "sent" in by, not "str" in by])
        and not strict
    ):
        if verbose:
            print(f"splited by digital (numbers)")
        return re.split(r"(\d+)", text)
    elif all([("digi" in by or "num" in by), "str" in by]) and not strict:
        if verbose:
            print(f"Splitting by (number strings)")
        pattern = re.compile(
            r"\b((?:one|two|three|four|five|six|seven|eight|nine|ten|eleven|twelve|thirteen|fourteen|fifteen|sixteen|seventeen|eighteen|nineteen|twenty|thirty|forty|fifty|sixty|seventy|eighty|ninety|hundred|thousand|million|billion|trillion|and|[\d,]+(?:\.\d+)?)(?:[-\s]?(?:one|two|three|four|five|six|seven|eight|nine|ten|eleven|twelve|thirteen|fourteen|fifteen|sixteen|seventeen|eighteen|nineteen|twenty|thirty|forty|fifty|sixty|seventy|eighty|ninety|hundred|thousand|million|billion|trillion|and|[\d,]+(?:\.\d+)?))*)\b",
            re.IGNORECASE,
        )
        return re.split(pattern, text)
    elif ("pun" in by) and not strict:
        if verbose:
            print(f"splited by 标点('.!?;')")
        return re.split(r"[.!?;]", text)
    elif ("\n" in by or "li" in by) and not strict:
        if verbose:
            print(f"splited by lines('\n')")
        return text.splitlines()
    elif ("cam" in by) and not strict:
        if verbose:
            print(f"splited by camel_case")
        return split_by_camel_case(text)
    elif ("word" in by) and not strict:
        if verbose:
            print(f"splited by word")
        return word_tokenize(text)
    elif ("sen" in by and not "num" in by) and not strict:
        if verbose:
            print(f"splited by sentence")
        return sent_tokenize(text)
    elif ("sen" in by and "num" in by) and not strict:
        return split_by_sent_num(text, **kws)
    elif ("cha" in by) and not strict:
        if verbose:
            print(f"splited by chracters")
        return list(text)
    elif ("up" in by or "cap" in by) and ("l" not in by) and not strict:
        if verbose:
            print(f"splited by upper case")
        return split_at_upper(text)
    elif ("u" in by and "l" in by) and not strict:
        if by.find("u") < by.find("l"):
            if verbose:
                print(f"splited by upper followed by lower case")
            return split_at_upper_fl_lower(text)
        else:
            if verbose:
                print(f"splited by lower followed by upper case")
            return split_at_lower_fl_upper(text)
    elif ("start" in by or "head" in by) and not strict:
        if verbose:
            print(f"splited by lookahead")
        return split_by_regex_lookahead(text, **kws)
    elif ("end" in by or "tail" in by) and not strict:
        if verbose:
            print(f"splited by endings")
        return split_by_regex_end(text, **kws)
    elif ("other" in by or "non_alp" in by) and not strict:
        if verbose:
            print(f"splited by non_alphanumeric")
        return split_by_consecutive_non_alphanumeric(text)
    elif ("len" in by) and not strict:
        if verbose:
            print(f"splited by fixed length")
        return split_by_fixed_length_chunks(text, **kws)
    elif ("re" in by or "cus" in by or "cos" in by) and not strict:
        if verbose:
            print(f"splited by customed, re; => {by}")
        return reg_split(text, **kws)
    elif ("lang" in by or "eng" in by) and not strict:
        return split_non_ascii(text)
    else:
        return split_general(text, by, verbose=verbose, **kws)


def pdf2img(dir_pdf, dir_save=None, page=None, kind="png", verbose=True, **kws):
    df_dir_img_single_page = pd.DataFrame()
    dir_single_page = []
    if verbose:
        pp(pdfinfo_from_path(dir_pdf))
    if isinstance(page, tuple) and page:
        page = list(page)
    if isinstance(page, int):
        page = [page]
    if page is None:
        page = [pdfinfo_from_path(dir_pdf)["Pages"]]
    if len(page) == 1 and page != [pdfinfo_from_path(dir_pdf)["Pages"]]:
        page = [page[0], page[0]]
    else:
        page = [1, page[0]]
    print(page)
    pages = convert_from_path(dir_pdf, first_page=page[0], last_page=page[-1], **kws)
    if dir_save is None:
        dir_save = mkdir(dirname(dir_pdf), basename(dir_pdf).split(".")[0] + "_img")
    for i, page in enumerate(pages):
        if verbose:
            print(f"processing page: {i+1}")
        if i < 9:
            dir_img_each_page = dir_save + f"page_0{i+1}.png"
        else:
            dir_img_each_page = dir_save + f"page_{i+1}.png"
        dir_single_page.append(dir_img_each_page)
        page.save(dir_img_each_page, kind.upper())
    df_dir_img_single_page["fpath"] = dir_single_page
    return df_dir_img_single_page


# dir_pdf = "/Users/macjianfeng/Dropbox/github/python/240308_Python Data Science Handbook.pdf"
# df_page = pdf2img(dir_pdf, page=[1, 5],dpi=300)
def get_encoding(fpath, alternative_encodings=None, verbose=False):
    """
    Attempt to determine the encoding of a file by trying multiple encodings.

    Parameters:
    fpath (str): The path to the file.
    alternative_encodings (list): List of encodings to try. If None, uses a default list.
    verbose (bool): If True, print detailed information about each attempted encoding.

    Returns:
    str: The encoding that successfully read the file, or None if no encoding worked.
    """
    if alternative_encodings is None:
        alternative_encodings = [
            "utf-8",
            "latin1",
            "windows-1252",
            "iso-8859-1",
            "iso-8859-2",
            "iso-8859-3",
            "iso-8859-4",
            "iso-8859-5",
            "iso-8859-6",
            "iso-8859-7",
            "iso-8859-8",
            "iso-8859-9",
            "windows-1250",
            "windows-1251",
            "windows-1253",
            "windows-1254",
            "windows-1255",
            "windows-1256",
            "windows-1257",
            "windows-1258",
            "big5",
            "gb18030",
            "shift_jis",
            "euc_jp",
            "koi8_r",
            "mac_roman",
            "mac_central_europe",
            "mac_greek",
            "mac_cyrillic",
            "mac_arabic",
            "mac_hebrew",
        ]

    if not os.path.isfile(fpath):
        raise FileNotFoundError(f"The file {fpath} does not exist.")

    for enc in alternative_encodings:
        try:
            with open(fpath, mode="r", encoding=enc) as file:
                file.read()  # Try to read the file
            if verbose:
                print(f"Successfully detected encoding: {enc}")
            return enc
        except UnicodeDecodeError:
            if verbose:
                print(f"Failed to decode with encoding: {enc}")
            continue

    # If no encoding worked
    print("No suitable encoding found.")
    return None


def fload(fpath, kind=None, **kwargs):
    """
    Load content from a file with specified file type.
    Parameters:
        fpath (str): The file path from which content will be loaded.
        kind (str): The file type to load. Supported options: 'docx', 'txt', 'md', 'html', 'json', 'yaml', 'xml', 'csv', 'xlsx', 'pdf'.
        **kwargs: Additional parameters for 'csv' and 'xlsx' file types.
    Returns:
        content: The content loaded from the file.
    """

    def load_txt_md(fpath):
        with open(fpath, "r") as file:
            content = file.read()
        return content

    def load_html(fpath):
        with open(fpath, "r") as file:
            content = file.read()
        return content

    def load_json(fpath):
        with open(fpath, "r") as file:
            content = json.load(file)
        return content

    def load_yaml(fpath):
        with open(fpath, "r") as file:
            content = yaml.safe_load(file)
        return content

    def load_xml(fpath):
        tree = etree.parse(fpath)
        root = tree.getroot()
        return etree.tostring(root, pretty_print=True).decode()

    def load_csv(fpath, engine="pyarrow", **kwargs):
        print(f"engine={engine}")
        df = pd.read_csv(fpath, engine=engine, **kwargs)
        return df

    def load_xlsx(fpath, **kwargs):
        df = pd.read_excel(fpath, **kwargs)
        return df

    def load_ipynb(fpath, **kwargs):
        as_version = kwargs.get("as_version", 4)
        with open(fpath, "r") as file:
            nb = nbformat.read(file, as_version=as_version)
            md_exporter = MarkdownExporter()
            md_body, _ = md_exporter.from_notebook_node(nb)
        return md_body

    def load_pdf(fpath, page="all", verbose=False, **kwargs):
        """
        Parameters:
        fpath: The path to the PDF file to be loaded.
        page (optional):
            Specifies which page or pages to extract text from. By default, it's set to "all", which means text from all
            pages will be returned. It can also be an integer to specify a single page number or a list of integers to
            specify multiple pages.
        verbose (optional):
            If True, prints the total number of pages processed.
        Functionality:
        It initializes an empty dictionary text_dict to store page numbers as keys and their corresponding text as values.
        It iterates through each page of the PDF file using a for loop.
        For each page, it extracts the text using PyPDF2's extract_text() method and stores it in text_dict with the page number incremented by 1 as the key.
        If the page parameter is an integer, it converts it into a list containing that single page number to ensure consistency in handling.
        If the page parameter is a NumPy array, it converts it to a list using the tolist() method to ensure compatibility with list operations.
        If verbose is True, it prints the total number of pages processed.
        If page is a list, it combines the text of the specified pages into a single string combined_text and returns it.
        If page is set to "all", it returns the entire text_dict containing text of all pages.
        If page is an integer, it returns the text of the specified page number.
        If the specified page is not found, it returns the string "Page is not found".
        """
        text_dict = {}
        with open(fpath, "rb") as file:
            pdf_reader = PdfReader(file)
            num_pages = len(pdf_reader.pages)
            for page_num in range(num_pages):
                if verbose:
                    print(f"processing page {page_num}")
                page_ = pdf_reader.pages[page_num]
                text_dict[page_num + 1] = page_.extract_text()
        if isinstance(page, int):
            page = [page]
        elif isinstance(page, np.ndarray):
            page = page.tolist()
        if verbose:
            print(f"total pages: {page_num}")
        if isinstance(page, list):
            combined_text = ""
            for page_num in page:
                combined_text += text_dict.get(page_num, "")
            return combined_text
        elif "all" in page.lower():
            combined_text = ""
            for i in text_dict.values():
                combined_text += i
            return combined_text
        else:
            return text_dict.get(int(page), "Page is not found")

    def load_docx(fpath):
        doc = Document(fpath)
        content = [para.text for para in doc.paragraphs]
        return content

    if kind is None:
        _, kind = os.path.splitext(fpath)
        kind = kind.lower()

    kind = kind.lstrip(".").lower()
    img_types = [
        "bmp",
        "eps",
        "gif",
        "icns",
        "ico",
        "im",
        "jpg",
        "jpeg",
        "jpeg2000",
        "msp",
        "pcx",
        "png",
        "ppm",
        "sgi",
        "spider",
        "tga",
        "tiff",
        "tif",
        "webp",
        "json",
    ]
    doc_types = [
        "docx",
        "txt",
        "md",
        "html",
        "json",
        "yaml",
        "xml",
        "csv",
        "xlsx",
        "pdf",
        "ipynb",
    ]
    supported_types = [*doc_types, *img_types]
    if kind not in supported_types:
        raise ValueError(
            f"Error:\n{kind} is not in the supported list {supported_types}"
        )
    if kind == "docx":
        return load_docx(fpath)
    elif kind == "txt" or kind == "md":
        return load_txt_md(fpath)
    elif kind == "html":
        return load_html(fpath)
    elif kind == "json":
        return load_json(fpath)
    elif kind == "yaml":
        return load_yaml(fpath)
    elif kind == "xml":
        return load_xml(fpath)
    elif kind == "csv":
        return load_csv(fpath, **kwargs)
    elif kind == "xlsx":
        return load_xlsx(fpath, **kwargs)
    elif kind == "ipynb":
        return load_ipynb(fpath, **kwargs)
    elif kind == "pdf":
        # print('usage:load_pdf(fpath, page="all", verbose=False)')
        return load_pdf(fpath, **kwargs)
    elif kind.lower() in img_types:
        print(f'Image ".{kind}" is loaded.')
        return load_img(fpath)
    else:
        raise ValueError(
            f"Error:\n{kind} is not in the supported list {supported_types}"
        )


# Example usage
# txt_content = fload('sample.txt')
# md_content = fload('sample.md')
# html_content = fload('sample.html')
# json_content = fload('sample.json')
# yaml_content = fload('sample.yaml')
# xml_content = fload('sample.xml')
# csv_content = fload('sample.csv')
# xlsx_content = fload('sample.xlsx')
# docx_content = fload('sample.docx')


def fupdate(fpath, content=None):
    """
    Update a file by adding new content at the top and moving the old content to the bottom.
    Parameters
    ----------
    fpath : str
        The file path where the content should be updated.
    content : str, optional
        The new content to add at the top of the file. If not provided, the function will not add any new content.
    Notes
    -----
    - If the file at `fpath` does not exist, it will be created.
    - The new content will be added at the top, followed by the old content of the file.
    """
    content = content or ""
    if os.path.exists(fpath):
        with open(fpath, "r") as file:
            old_content = file.read()
    else:
        old_content = ""

    with open(fpath, "w") as file:
        file.write(content)
        file.write(old_content)


def fappend(fpath, content=None):
    """
    append new content at the end.
    """
    content = content or ""
    if os.path.exists(fpath):
        with open(fpath, "r") as file:
            old_content = file.read()
    else:
        old_content = ""

    with open(fpath, "w") as file:
        file.write(old_content)
        file.write(content)


def fsave(
    fpath,
    content,
    mode="w",
    how="overwrite",
    kind=None,
    font_name="Times",
    font_size=10,
    spacing=6,
    **kwargs,
):
    """
    Save content into a file with specified file type and formatting.
    Parameters:
        fpath (str): The file path where content will be saved.
        content (list of str or dict): The content to be saved, where each string represents a paragraph or a dictionary for tabular data.
        kind (str): The file type to save. Supported options: 'docx', 'txt', 'md', 'html', 'pdf', 'csv', 'xlsx', 'json', 'xml', 'yaml'.
        font_name (str): The font name for text formatting (only applicable for 'docx', 'html', and 'pdf').
        font_size (int): The font size for text formatting (only applicable for 'docx', 'html', and 'pdf').
        spacing (int): The space after each paragraph (only applicable for 'docx').
        **kwargs: Additional parameters for 'csv', 'xlsx', 'json', 'yaml' file types.
    Returns:
        None
    """

    def save_content(fpath, content, mode=mode, how="overwrite"):
        if "wri" in how.lower():
            with open(fpath, mode, encoding="utf-8") as file:
                file.write(content)
        elif "upd" in how.lower():
            fupdate(fpath, content=content)
        elif "app" in how.lower():
            fappend(fpath, content=content)

    def save_docx(fpath, content, font_name, font_size, spacing):
        if isinstance(content, str):
            content = content.split(". ")
        doc = docx.Document()
        for i, paragraph_text in enumerate(content):
            paragraph = doc.add_paragraph()
            run = paragraph.add_run(paragraph_text)
            font = run.font
            font.name = font_name
            font.size = docx.shared.Pt(font_size)
            if i != len(content) - 1:  # Add spacing for all but the last paragraph
                paragraph.space_after = docx.shared.Pt(spacing)
        doc.save(fpath)

    def save_txt_md(fpath, content, sep="\n", mode="w"):
        # Ensure content is a single string
        if isinstance(content, list):
            content = sep.join(content)
        save_content(fpath, sep.join(content), mode)

    def save_html(fpath, content, font_name, font_size, mode="w"):
        html_content = "<html><body>"
        for paragraph_text in content:
            html_content += f'<p style="font-family:{font_name}; font-size:{font_size}px;">{paragraph_text}</p>'
        html_content += "</body></html>"
        save_content(fpath, html_content, mode)

    def save_pdf(fpath, content, font_name, font_size):
        pdf = FPDF()
        pdf.add_page()
        # pdf.add_font('Arial','',r'/System/Library/Fonts/Supplemental/Arial.ttf',uni=True)
        pdf.set_font(font_name, "", font_size)
        for paragraph_text in content:
            pdf.multi_cell(0, 10, paragraph_text)
            pdf.ln(h="")
        pdf.output(fpath, "F")

    def save_csv(fpath, data, **kwargs):
        df = pd.DataFrame(data)
        df.to_csv(fpath, **kwargs)

    def save_xlsx(fpath, data, **kwargs):
        df = pd.DataFrame(data)
        df.to_excel(fpath, **kwargs)

    def save_ipynb(fpath, data, **kwargs):
        # Split the content by code fences to distinguish between code and markdown
        parts = data.split("```")
        cells = []

        for i, part in enumerate(parts):
            if i % 2 == 0:
                # Even index: markdown content
                cells.append(nbformat.v4.new_markdown_cell(part.strip()))
            else:
                # Odd index: code content
                cells.append(nbformat.v4.new_code_cell(part.strip()))
        # Create a new notebook
        nb = nbformat.v4.new_notebook()
        nb["cells"] = cells
        # Write the notebook to a file
        with open(fpath, "w", encoding="utf-8") as ipynb_file:
            nbformat.write(nb, ipynb_file)

    # def save_json(fpath, data, **kwargs):
    #     with open(fpath, "w") as file:
    #         json.dump(data, file, **kwargs)

    def save_json(fpath_fname, var_dict_or_df):
        with open(fpath_fname, "w") as f_json:
            if isinstance(var_dict_or_df, pd.DataFrame):
                var_dict_or_df = var_dict_or_df.to_dict(orient="dict")
            if isinstance(var_dict_or_df, dict):
                for key, value in var_dict_or_df.items():
                    if isinstance(value, np.ndarray):
                        var_dict_or_df[key] = value.tolist()
            # Save the dictionary or list of dictionaries to a JSON file
            json.dump(var_dict_or_df, f_json, indent=4)

    # # Example usage:
    # sets = {"title": "mse_path_ MSE"}
    # jsonsave("/.json", sets)
    # # setss = jsonload("/.json")

    def save_yaml(fpath, data, **kwargs):
        with open(fpath, "w") as file:
            yaml.dump(data, file, **kwargs)

    def save_xml(fpath, data):
        root = etree.Element("root")
        if isinstance(data, dict):
            for key, val in data.items():
                child = etree.SubElement(root, key)
                child.text = str(val)
        else:
            raise ValueError("XML saving only supports dictionary data")
        tree = etree.ElementTree(root)
        tree.write(fpath, pretty_print=True, xml_declaration=True, encoding="UTF-8")

    if kind is None:
        _, kind = os.path.splitext(fpath)
        kind = kind.lower()

    kind = kind.lstrip(".").lower()

    if kind not in [
        "docx",
        "txt",
        "md",
        "html",
        "pdf",
        "csv",
        "xlsx",
        "json",
        "xml",
        "yaml",
        "ipynb",
    ]:
        print(
            f"Warning:\n{kind} is not in the supported list ['docx', 'txt', 'md', 'html', 'pdf', 'csv', 'xlsx', 'json', 'xml', 'yaml']"
        )

    if kind == "docx" or kind == "doc":
        save_docx(fpath, content, font_name, font_size, spacing)
    elif kind == "txt":
        save_txt_md(fpath, content, sep="", mode=mode)
    elif kind == "md":
        save_txt_md(fpath, content, sep="", mode=mode)
    elif kind == "html":
        save_html(fpath, content, font_name, font_size)
    elif kind == "pdf":
        save_pdf(fpath, content, font_name, font_size)
    elif kind == "csv":
        save_csv(fpath, content, **kwargs)
    elif kind == "xlsx":
        save_xlsx(fpath, content, **kwargs)
    elif kind == "json":
        save_json(fpath, content)
    elif kind == "xml":
        save_xml(fpath, content)
    elif kind == "yaml":
        save_yaml(fpath, content, **kwargs)
    elif kind == "ipynb":
        save_ipynb(fpath, content, **kwargs)
    else:
        try:
            netfinder.downloader(url=content, dir_save=dirname(fpath), kind=kind)
        except:
            print(
                f"Error:\n{kind} is not in the supported list ['docx', 'txt', 'md', 'html', 'pdf', 'csv', 'xlsx', 'json', 'xml', 'yaml']"
            )


# # Example usage
# text_content = ["Hello, this is a sample text file.", "This is the second paragraph."]
# tabular_content = {"Name": ["Alice", "Bob"], "Age": [24, 30]}
# json_content = {"name": "Alice", "age": 24}
# yaml_content = {"Name": "Alice", "Age": 24}
# xml_content = {"Name": "Alice", "Age": 24}
# dir_save = "/Users/macjianfeng/Dropbox/Downloads/"
# fsave(dir_save + "sample.txt", text_content)
# fsave(dir_save + "sample.md", text_content)
# fsave(dir_save + "sample.html", text_content)
# fsave(dir_save + "sample.pdf", text_content)
# fsave(dir_save + "sample.docx", text_content)
# fsave(dir_save + "sample.csv", tabular_content, index=False)
# fsave(dir_save + "sample.xlsx", tabular_content, sheet_name="Sheet1", index=False)
# fsave(dir_save + "sample.json", json_content, indent=4)
# fsave(dir_save + "sample.yaml", yaml_content)
# fsave(dir_save + "sample.xml", xml_content)


def addpath(fpath):
    sys.path.insert(0, dir)


def dirname(fpath):
    """
    dirname: Extracting Directory Name from a File Path
    Args:
        fpath (str): the file or directory path
    Returns:
        str: directory, without filename
    """
    dirname_ = os.path.dirname(fpath)
    if not dirname_.endswith("/"):
        dirname_ = dirname_ + "/"
    return dirname_


def dir_name(fpath):  # same as "dirname"
    return dirname(fpath)


def basename(fpath):
    """
    basename: # Output: file.txt
    Args:
        fpath (str): the file or directory path
    Returns:
        str: # Output: file.txt
    """
    return os.path.basename(fpath)


def flist(fpath, contains="all"):
    all_files = [
        os.path.join(fpath, f)
        for f in os.listdir(fpath)
        if os.path.isfile(os.path.join(fpath, f))
    ]
    if isinstance(contains, list):
        filt_files = []
        for filter_ in contains:
            filt_files.extend(flist(fpath, filter_))
        return filt_files
    else:
        if "all" in contains.lower():
            return all_files
        else:
            filt_files = [f for f in all_files if isa(f, contains)]
            return filt_files


def sort_kind(df, by="name", ascending=True):
    if df[by].dtype == "object":  # Check if the column contains string values
        if ascending:
            sorted_index = df[by].str.lower().argsort()
        else:
            sorted_index = df[by].str.lower().argsort()[::-1]
    else:
        if ascending:
            sorted_index = df[by].argsort()
        else:
            sorted_index = df[by].argsort()[::-1]
    sorted_df = df.iloc[sorted_index].reset_index(drop=True)
    return sorted_df


def isa(*args, **kwargs):
    """
    fpath, contains='img'
    containss file paths based on the specified contains.
    Args:
        fpath (str): Path to the file.
        contains (str): contains of file to contains. Default is 'img' for images. Other options include 'doc' for documents,
                    'zip' for ZIP archives, and 'other' for other types of files.
    Returns:
        bool: True if the file matches the contains, False otherwise.
    """
    for arg in args:
        if isinstance(arg, str):
            if "/" in arg or "\\" in arg:
                fpath = arg
            else:
                contains = arg
    if "img" in contains.lower() or "image" in contains.lower():
        return is_image(fpath)
    elif "doc" in contains.lower():
        return is_document(fpath)
    elif "zip" in contains.lower():
        return is_zip(fpath)
    elif "dir" in contains.lower() or (
        "f" in contains.lower() and "d" in contains.lower()
    ):
        return os.path.isdir(fpath)
    elif "fi" in contains.lower():  # file
        return os.path.isfile(fpath)
    elif "num" in contains.lower():  # file
        return os.path.isfile(fpath)
    elif "text" in contains.lower() or "txt" in contains.lower():  # file
        return is_text(fpath)
    elif "color" in contains.lower():  # file
        return is_str_color(fpath)
    else:
        print(f"{contains} was not set up correctly")
        return False


def listdir(
    rootdir,
    kind="folder",
    sort_by="name",
    ascending=True,
    contains=None,
    orient="list",
    output="df",  # 'list','dict','records','index','series'
):
    if not kind.startswith("."):
        kind = "." + kind

    if os.path.isdir(rootdir):
        ls = os.listdir(rootdir)
        fd = [".fd", ".fld", ".fol", ".fd", ".folder"]
        i = 0
        f = {
            "name": [],
            "length": [],
            "path": [],
            "created_time": [],
            "modified_time": [],
            "last_open_time": [],
            "size": [],
            "fname": [],
            "fpath": [],
        }
        for item in ls:
            item_path = os.path.join(rootdir, item)
            if item.startswith("."):
                continue
            filename, file_extension = os.path.splitext(item)
            is_folder = kind.lower() in fd and os.path.isdir(item_path)
            is_file = kind.lower() in file_extension.lower() and (
                os.path.isfile(item_path)
            )
            if kind in [".doc", ".img", ".zip"]:  # 选择大的类别
                if kind != ".folder" and not isa(item_path, kind):
                    continue
            elif kind in [".all"]:
                return flist(fpath, contains=contains)
            else:  # 精确到文件的后缀
                if not is_folder and not is_file:
                    continue
            f["name"].append(filename)
            f["length"].append(len(filename))
            f["path"].append(os.path.join(os.path.dirname(item_path), item))
            fpath = os.path.join(os.path.dirname(item_path), item)
            f["size"].append(round(os.path.getsize(fpath) / 1024 / 1024, 3))
            f["created_time"].append(
                pd.to_datetime(os.path.getctime(item_path), unit="s")
            )
            f["modified_time"].append(
                pd.to_datetime(os.path.getmtime(item_path), unit="s")
            )
            f["last_open_time"].append(
                pd.to_datetime(os.path.getatime(item_path), unit="s")
            )
            f["fname"].append(filename)  # will be removed
            f["fpath"].append(fpath)  # will be removed
            i += 1

        f["num"] = i
        f["rootdir"] = rootdir
        f["os"] = os.uname().machine
    else:
        raise FileNotFoundError(
            'The directory "{}" does NOT exist. Please check the directory "rootdir".'.format(
                rootdir
            )
        )

    f = pd.DataFrame(f)

    if contains is not None:
        f = f[f["name"].str.contains(contains, case=False)]

    if "nam" in sort_by.lower():
        f = sort_kind(f, by="name", ascending=ascending)
    elif "crea" in sort_by.lower():
        f = sort_kind(f, by="created_time", ascending=ascending)
    elif "modi" in sort_by.lower():
        f = sort_kind(f, by="modified_time", ascending=ascending)
    elif "s" in sort_by.lower() and "z" in sort_by.lower():
        f = sort_kind(f, by="size", ascending=ascending)

    if "df" in output:
        return f
    else:
        if "l" in orient.lower():  # list # default
            res_output = Box(f.to_dict(orient="list"))
            return res_output
        if "d" in orient.lower():  # dict
            return Box(f.to_dict(orient="dict"))
        if "r" in orient.lower():  # records
            return Box(f.to_dict(orient="records"))
        if "in" in orient.lower():  # records
            return Box(f.to_dict(orient="index"))
        if "se" in orient.lower():  # records
            return Box(f.to_dict(orient="series"))


# Example usage:
# result = listdir('your_root_directory')
# print(result)
# df=listdir("/", contains='sss',sort_by='name',ascending=False)
# print(df.fname.to_list(),"\n",df.fpath.to_list())
def list_func(lib_name, opt="call"):
    if opt == "call":
        funcs = [func for func in dir(lib_name) if callable(getattr(lib_name, func))]
    else:
        funcs = dir(lib_name)
    return funcs


def func_list(lib_name, opt="call"):
    return list_func(lib_name, opt=opt)


def mkdir(*args, **kwargs):
    """
    newfolder(pardir, chdir)
    Args:
        pardir (dir): parent dir
        chdir (str): children dir
        overwrite (bool): overwrite?
    Returns:
        mkdir, giving a option if exists_ok or not
    """
    overwrite = kwargs.get("overwrite", False)
    for arg in args:
        if isinstance(arg, (str, list)):
            if "/" in arg or "\\" in arg:
                pardir = arg
                print(f"pardir{pardir}")
            else:
                chdir = arg
                print(f"chdir{chdir}")
        elif isinstance(arg, bool):
            overwrite = arg
            print(overwrite)
        else:
            print(f"{arg}Error: not support a {type(arg)} type")
    rootdir = []
    # Convert string to list
    if isinstance(chdir, str):
        chdir = [chdir]
    # Subfoldername should be unique
    chdir = list(set(chdir))
    if isinstance(pardir, str):  # Dir_parents should be 'str' type
        pardir = os.path.normpath(pardir)
    # Get the slash type: "/" or "\"
    stype = "/" if "/" in pardir else "\\"
    # Check if the parent directory exists and is a directory path
    if os.path.isdir(pardir):
        os.chdir(pardir)  # Set current path
        # Check if subdirectories are not empty
        if chdir:
            chdir.sort()
            # Create multiple subdirectories at once
            for folder in chdir:
                # Check if the subfolder already exists
                child_tmp = os.path.join(pardir, folder)
                if not os.path.isdir(child_tmp):
                    os.mkdir("./" + folder)
                    print(f"\n {folder} was created successfully!\n")
                else:
                    if overwrite:
                        shutil.rmtree(child_tmp)
                        os.mkdir("./" + folder)
                        print(f"\n {folder} overwrite! \n")
                    else:
                        print(f"\n {folder} already exists! \n")
                rootdir.append(child_tmp + stype)  # Note down
        else:
            print("\nWarning: Dir_child doesn't exist\n")
    else:
        print("\nWarning: Dir_parent is not a directory path\n")
    # Dir is the main output, if only one dir, then str type is inconvenient
    if len(rootdir) == 1:
        rootdir = rootdir[0]
    return rootdir


def figsave(*args, dpi=300):
    dir_save = None
    fname = None
    for arg in args:
        if isinstance(arg, str):
            if "/" in arg or "\\" in arg:
                dir_save = arg
            elif "/" not in arg and "\\" not in arg:
                fname = arg
    # Backup original values
    if "/" in dir_save:
        if dir_save[-1] != "/":
            dir_save = dir_save + "/"
    elif "\\" in dir_save:
        if dir_save[-1] != "\\":
            dir_save = dir_save + "\\"
    else:
        raise ValueError("Check the Path of dir_save Directory")
    ftype = fname.split(".")[-1]
    if len(fname.split(".")) == 1:
        ftype = "nofmt"
        fname = dir_save + fname + "." + ftype
    else:
        fname = dir_save + fname
    # Save figure based on file type
    if ftype.lower() == "eps":
        plt.savefig(fname, format="eps", bbox_inches="tight")
        plt.savefig(
            fname.replace(".eps", ".pdf"), format="pdf", bbox_inches="tight", dpi=dpi
        )
    elif ftype.lower() == "nofmt":  # default: both "tif" and "pdf"
        fname_corr = fname.replace("nofmt", "pdf")
        plt.savefig(fname_corr, format="pdf", bbox_inches="tight", dpi=dpi)
        fname = fname.replace("nofmt", "tif")
        plt.savefig(fname, format="tiff", dpi=dpi, bbox_inches="tight")
        print(f"default saving filetype: both 'tif' and 'pdf")
    elif ftype.lower() == "pdf":
        plt.savefig(fname, format="pdf", bbox_inches="tight", dpi=dpi)
    elif ftype.lower() in ["jpg", "jpeg"]:
        plt.savefig(fname, format="jpeg", dpi=dpi, bbox_inches="tight")
    elif ftype.lower() == "png":
        plt.savefig(fname, format="png", dpi=dpi, bbox_inches="tight", transparent=True)
    elif ftype.lower() in ["tiff", "tif"]:
        plt.savefig(fname, format="tiff", dpi=dpi, bbox_inches="tight")
    elif ftype.lower() == "emf":
        plt.savefig(fname, format="emf", dpi=dpi, bbox_inches="tight")
    elif ftype.lower() == "fig":
        plt.savefig(fname, format="pdf", bbox_inches="tight", dpi=dpi)
    print(f"\nSaved @: dpi={dpi}\n{fname}")


def is_str_color(s):
    # Regular expression pattern for hexadecimal color codes
    color_code_pattern = r"^#([A-Fa-f0-9]{6}|[A-Fa-f0-9]{8})$"
    return re.match(color_code_pattern, s) is not None


def is_num(s):
    """
    Check if a string can be converted to a number (int or float).
    Parameters:
    - s (str): The string to check.
    Returns:
    - bool: True if the string can be converted to a number, False otherwise.
    """
    try:
        float(s)  # Try converting the string to a float
        return True
    except ValueError:
        return False


def isnum(s):
    return is_num(s)


def is_image(fpath):
    mime_type, _ = mimetypes.guess_type(fpath)
    if mime_type and mime_type.startswith("image"):
        return True
    else:
        return False


def is_document(fpath):
    mime_type, _ = mimetypes.guess_type(fpath)
    if mime_type and (
        mime_type.startswith("text/")
        or mime_type == "application/pdf"
        or mime_type == "application/msword"
        or mime_type
        == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        or mime_type == "application/vnd.ms-excel"
        or mime_type
        == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        or mime_type == "application/vnd.ms-powerpoint"
        or mime_type
        == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    ):
        return True
    else:
        return False


def is_zip(fpath):
    mime_type, _ = mimetypes.guess_type(fpath)
    if mime_type == "application/zip":
        return True
    else:
        return False


def adjust_spines(ax=None, spines=["left", "bottom"], distance=2):
    if ax is None:
        ax = plt.gca()
    for loc, spine in ax.spines.items():
        if loc in spines:
            spine.set_position(("outward", distance))  # outward by 2 points
            # spine.set_smart_bounds(True)
        else:
            spine.set_color("none")  # don't draw spine
    # turn off ticks where there is no spine
    if "left" in spines:
        ax.yaxis.set_ticks_position("left")
    else:
        ax.yaxis.set_ticks([])
    if "bottom" in spines:
        ax.xaxis.set_ticks_position("bottom")
    else:
        # no xaxis ticks
        ax.xaxis.set_ticks([])


# And then plot the data:


def add_colorbar(im, width=None, pad=None, **kwargs):
    # usage: add_colorbar(im, width=0.01, pad=0.005, label="PSD (dB)", shrink=0.8)
    l, b, w, h = im.axes.get_position().bounds  # get boundaries
    width = width or 0.1 * w  # get width of the colorbar
    pad = pad or width  # get pad between im and cbar
    fig = im.axes.figure  # get figure of image
    cax = fig.add_axes([l + w + pad, b, width, h])  # define cbar Axes
    return fig.colorbar(im, cax=cax, **kwargs)  # draw cbar


# =============================================================================
# # for plot figures: setting rcParams
# usage: set_pub()
# or by setting sns.set_theme...see below:
# sns.set_theme(style="ticks", rc=params)      # 白色无刻度线，有坐标轴标度
# # sns.set_theme(style="whitegrid", rc=params)# 白色＋刻度线，无坐标轴标度
# # sns.set_theme(style="white", rc=params)    # 白色无刻度线，无坐标轴标度
# # sns.set_theme(style="dark", rc=params)     # 深色无刻度线，无坐标轴标度
# =============================================================================


def list2str(x_str):
    s = "".join(str(x) for x in x_str)
    return s


def str2list(str_):
    l = []
    [l.append(x) for x in str_]
    return l


def load_img(fpath):
    """
    Load an image from the specified file path.
    Args:
        fpath (str): The file path to the image.
    Returns:
        PIL.Image: The loaded image.
    Raises:
        FileNotFoundError: If the specified file is not found.
        OSError: If the specified file cannot be opened or is not a valid image file.
    """
    try:
        img = Image.open(fpath)
        return img
    except FileNotFoundError:
        raise FileNotFoundError(f"The file '{fpath}' was not found.")
    except OSError:
        raise OSError(f"Unable to open file '{fpath}' or it is not a valid image file.")


def apply_filter(img, *args):
    # def apply_filter(img, filter_name, filter_value=None):
    """
    Apply the specified filter to the image.
    Args:
        img (PIL.Image): The input image.
        filter_name (str): The name of the filter to apply.
        **kwargs: Additional parameters specific to the filter.
    Returns:
        PIL.Image: The filtered image.
    """

    def correct_filter_name(filter_name):
        if "bl" in filter_name.lower() and "box" not in filter_name.lower():
            return "BLUR"
        elif "cont" in filter_name.lower():
            return "Contour"
        elif "det" in filter_name.lower():
            return "Detail"
        elif (
            "edg" in filter_name.lower()
            and "mo" not in filter_name.lower()
            and "f" not in filter_name.lower()
        ):
            return "EDGE_ENHANCE"
        elif "edg" in filter_name.lower() and "mo" in filter_name.lower():
            return "EDGE_ENHANCE_MORE"
        elif "emb" in filter_name.lower():
            return "EMBOSS"
        elif "edg" in filter_name.lower() and "f" in filter_name.lower():
            return "FIND_EDGES"
        elif "sh" in filter_name.lower() and "mo" not in filter_name.lower():
            return "SHARPEN"
        elif "sm" in filter_name.lower() and "mo" not in filter_name.lower():
            return "SMOOTH"
        elif "sm" in filter_name.lower() and "mo" in filter_name.lower():
            return "SMOOTH_MORE"
        elif "min" in filter_name.lower():
            return "MIN_FILTER"
        elif "max" in filter_name.lower():
            return "MAX_FILTER"
        elif "mod" in filter_name.lower():
            return "MODE_FILTER"
        elif "mul" in filter_name.lower():
            return "MULTIBAND_FILTER"
        elif "gau" in filter_name.lower():
            return "GAUSSIAN_BLUR"
        elif "box" in filter_name.lower():
            return "BOX_BLUR"
        elif "med" in filter_name.lower():
            return "MEDIAN_FILTER"
        else:
            supported_filters = [
                "BLUR",
                "CONTOUR",
                "DETAIL",
                "EDGE_ENHANCE",
                "EDGE_ENHANCE_MORE",
                "EMBOSS",
                "FIND_EDGES",
                "SHARPEN",
                "SMOOTH",
                "SMOOTH_MORE",
                "MIN_FILTER",
                "MAX_FILTER",
                "MODE_FILTER",
                "MULTIBAND_FILTER",
                "GAUSSIAN_BLUR",
                "BOX_BLUR",
                "MEDIAN_FILTER",
            ]
            raise ValueError(
                f"Unsupported filter: {filter_name}, should be one of: {supported_filters}"
            )

    for arg in args:
        if isinstance(arg, str):
            filter_name = arg
            filter_name = correct_filter_name(filter_name)
        else:
            filter_value = arg
    filter_name = filter_name.upper()  # Ensure filter name is uppercase

    # Supported filters
    supported_filters = {
        "BLUR": ImageFilter.BLUR,
        "CONTOUR": ImageFilter.CONTOUR,
        "DETAIL": ImageFilter.DETAIL,
        "EDGE_ENHANCE": ImageFilter.EDGE_ENHANCE,
        "EDGE_ENHANCE_MORE": ImageFilter.EDGE_ENHANCE_MORE,
        "EMBOSS": ImageFilter.EMBOSS,
        "FIND_EDGES": ImageFilter.FIND_EDGES,
        "SHARPEN": ImageFilter.SHARPEN,
        "SMOOTH": ImageFilter.SMOOTH,
        "SMOOTH_MORE": ImageFilter.SMOOTH_MORE,
        "MIN_FILTER": ImageFilter.MinFilter,
        "MAX_FILTER": ImageFilter.MaxFilter,
        "MODE_FILTER": ImageFilter.ModeFilter,
        "MULTIBAND_FILTER": ImageFilter.MultibandFilter,
        "GAUSSIAN_BLUR": ImageFilter.GaussianBlur,
        "BOX_BLUR": ImageFilter.BoxBlur,
        "MEDIAN_FILTER": ImageFilter.MedianFilter,
    }
    # Check if the filter name is supported
    if filter_name not in supported_filters:
        raise ValueError(
            f"Unsupported filter: {filter_name}, should be one of: {[i.lower() for i in supported_filters.keys()]}"
        )

    # Apply the filter
    if filter_name.upper() in [
        "BOX_BLUR",
        "GAUSSIAN_BLUR",
        "MEDIAN_FILTER",
        "MIN_FILTER",
        "MAX_FILTER",
        "MODE_FILTER",
    ]:
        radius = filter_value if filter_value is not None else 2
        return img.filter(supported_filters[filter_name](radius))
    elif filter_name in ["MULTIBAND_FILTER"]:
        bands = filter_value if filter_value is not None else None
        return img.filter(supported_filters[filter_name](bands))
    else:
        if filter_value is not None:
            print(
                f"{filter_name} doesn't require a value for {filter_value}, but it remains unaffected"
            )
        return img.filter(supported_filters[filter_name])


def imgsetss(
    img,
    sets=None,
    show=True,
    show_axis=False,
    size=None,
    dpi=100,
    figsize=None,
    auto=False,
    filter_kws=None,
):
    """
    Apply various enhancements and filters to an image using PIL's ImageEnhance and ImageFilter modules.

    Args:
        img (PIL.Image): The input image.
        sets (dict): A dictionary specifying the enhancements, filters, and their parameters.
        show (bool): Whether to display the enhanced image.
        show_axis (bool): Whether to display axes on the image plot.
        size (tuple): The size of the thumbnail, cover, contain, or fit operation.
        dpi (int): Dots per inch for the displayed image.
        figsize (tuple): The size of the figure for displaying the image.
        auto (bool): Whether to automatically enhance the image based on its characteristics.

    Returns:
        PIL.Image: The enhanced image.

    Supported enhancements and filters:
        - "sharpness": Adjusts the sharpness of the image. Values > 1 increase sharpness, while values < 1 decrease sharpness.
        - "contrast": Adjusts the contrast of the image. Values > 1 increase contrast, while values < 1 decrease contrast.
        - "brightness": Adjusts the brightness of the image. Values > 1 increase brightness, while values < 1 decrease brightness.
        - "color": Adjusts the color saturation of the image. Values > 1 increase saturation, while values < 1 decrease saturation.
        - "rotate": Rotates the image by the specified angle.
        - "crop" or "cut": Crops the image. The value should be a tuple specifying the crop box as (left, upper, right, lower).
        - "size": Resizes the image to the specified dimensions.
        - "thumbnail": Resizes the image to fit within the given size while preserving aspect ratio.
        - "cover": Resizes and crops the image to fill the specified size.
        - "contain": Resizes the image to fit within the specified size, adding borders if necessary.
        - "fit": Resizes and pads the image to fit within the specified size.
        - "filter": Applies various filters to the image (e.g., BLUR, CONTOUR, EDGE_ENHANCE).

    Note:
        The "color" and "enhance" enhancements are not implemented in this function.
    """
    supported_filters = [
        "BLUR",
        "CONTOUR",
        "DETAIL",
        "EDGE_ENHANCE",
        "EDGE_ENHANCE_MORE",
        "EMBOSS",
        "FIND_EDGES",
        "SHARPEN",
        "SMOOTH",
        "SMOOTH_MORE",
        "MIN_FILTER",
        "MAX_FILTER",
        "MODE_FILTER",
        "MULTIBAND_FILTER",
        "GAUSSIAN_BLUR",
        "BOX_BLUR",
        "MEDIAN_FILTER",
    ]
    print(
        "sets: a dict,'sharp:1.2','color','contrast:'auto' or 1.2','bright', 'crop: x_upperleft,y_upperleft, x_lowerright, y_lowerright','rotation','resize','rem or background'"
    )
    print(f"usage: filter_kws 'dict' below:")
    pp([str(i).lower() for i in supported_filters])
    print("\nlog:\n")

    def confirm_rembg_models(model_name):
        models_support = [
            "u2net",
            "u2netp",
            "u2net_human_seg",
            "u2net_cloth_seg",
            "silueta",
            "isnet-general-use",
            "isnet-anime",
            "sam",
        ]
        if model_name in models_support:
            print(f"model_name: {model_name}")
            return model_name
        else:
            print(
                f"{model_name} cannot be found, check the name:{models_support}, default('isnet-general-use') has been used"
            )
            return "isnet-general-use"

    def auto_enhance(img):
        """
        Automatically enhances the image based on its characteristics.
        Args:
            img (PIL.Image): The input image.
        Returns:
            dict: A dictionary containing the optimal enhancement values.
        """
        # Determine the bit depth based on the image mode
        if img.mode in ["1", "L", "P", "RGB", "YCbCr", "LAB", "HSV"]:
            # 8-bit depth per channel
            bit_depth = 8
        elif img.mode in ["RGBA", "CMYK"]:
            # 8-bit depth per channel + alpha (RGBA) or additional channels (CMYK)
            bit_depth = 8
        elif img.mode in ["I", "F"]:
            # 16-bit depth per channel (integer or floating-point)
            bit_depth = 16
        else:
            raise ValueError("Unsupported image mode")
        # Calculate the brightness and contrast for each channel
        num_channels = len(img.getbands())
        brightness_factors = []
        contrast_factors = []
        for channel in range(num_channels):
            channel_histogram = img.split()[channel].histogram()
            brightness = sum(i * w for i, w in enumerate(channel_histogram)) / sum(
                channel_histogram
            )
            channel_min, channel_max = img.split()[channel].getextrema()
            contrast = channel_max - channel_min
            # Adjust calculations based on bit depth
            normalization_factor = 2**bit_depth - 1  # Max value for the given bit depth
            brightness_factor = (
                1.0 + (brightness - normalization_factor / 2) / normalization_factor
            )
            contrast_factor = (
                1.0 + (contrast - normalization_factor / 2) / normalization_factor
            )
            brightness_factors.append(brightness_factor)
            contrast_factors.append(contrast_factor)
        # Calculate the average brightness and contrast factors across channels
        avg_brightness_factor = sum(brightness_factors) / num_channels
        avg_contrast_factor = sum(contrast_factors) / num_channels
        return {"brightness": avg_brightness_factor, "contrast": avg_contrast_factor}

    # Load image if input is a file path
    if isinstance(img, str):
        img = load_img(img)
    img_update = img.copy()
    # Auto-enhance image if requested
    if auto:
        auto_params = auto_enhance(img_update)
        sets.update(auto_params)
    if sets is None:
        sets = {}
    for k, value in sets.items():
        if "shar" in k.lower():
            enhancer = ImageEnhance.Sharpness(img_update)
            img_update = enhancer.enhance(value)
        elif "col" in k.lower() and "bg" not in k.lower():
            enhancer = ImageEnhance.Color(img_update)
            img_update = enhancer.enhance(value)
        elif "contr" in k.lower():
            if value and isinstance(value, (float, int)):
                enhancer = ImageEnhance.Contrast(img_update)
                img_update = enhancer.enhance(value)
            else:
                print("autocontrasted")
                img_update = ImageOps.autocontrast(img_update)
        elif "bri" in k.lower():
            enhancer = ImageEnhance.Brightness(img_update)
            img_update = enhancer.enhance(value)
        elif "cro" in k.lower() or "cut" in k.lower():
            img_update = img_update.crop(value)
        elif "rota" in k.lower():
            img_update = img_update.rotate(value)
        elif "si" in k.lower():
            img_update = img_update.resize(value)
        elif "thum" in k.lower():
            img_update.thumbnail(value)
        elif "cover" in k.lower():
            img_update = ImageOps.cover(img_update, size=value)
        elif "contain" in k.lower():
            img_update = ImageOps.contain(img_update, size=value)
        elif "fit" in k.lower():
            img_update = ImageOps.fit(img_update, size=value)
        elif "pad" in k.lower():
            img_update = ImageOps.pad(img_update, size=value)
        elif "rem" in k.lower() or "rm" in k.lower() or "back" in k.lower():
            if value and isinstance(value, (int, float, list)):
                print(
                    'example usage: {"rm":[alpha_matting_background_threshold(20),alpha_matting_foreground_threshold(270),alpha_matting_erode_sive(11)]}'
                )
                print("https://github.com/danielgatis/rembg/blob/main/USAGE.md")
                #     ###            Parameters:
                #         data (Union[bytes, PILImage, np.ndarray]): The input image data.
                #         alpha_matting (bool, optional): Flag indicating whether to use alpha matting. Defaults to False.
                #         alpha_matting_foreground_threshold (int, optional): Foreground threshold for alpha matting. Defaults to 240.
                #         alpha_matting_background_threshold (int, optional): Background threshold for alpha matting. Defaults to 10.
                #         alpha_matting_erode_size (int, optional): Erosion size for alpha matting. Defaults to 10.
                #         session (Optional[BaseSession], optional): A session object for the 'u2net' model. Defaults to None.
                #         only_mask (bool, optional): Flag indicating whether to return only the binary masks. Defaults to False.
                #         post_process_mask (bool, optional): Flag indicating whether to post-process the masks. Defaults to False.
                #         bgcolor (Optional[Tuple[int, int, int, int]], optional): Background color for the cutout image. Defaults to None.
                #  ###
                if isinstance(value, int):
                    value = [value]
                if len(value) < 2:
                    img_update = remove(
                        img_update,
                        alpha_matting=True,
                        alpha_matting_background_threshold=value,
                    )
                elif 2 <= len(value) < 3:
                    img_update = remove(
                        img_update,
                        alpha_matting=True,
                        alpha_matting_background_threshold=value[0],
                        alpha_matting_foreground_threshold=value[1],
                    )
                elif 3 <= len(value) < 4:
                    img_update = remove(
                        img_update,
                        alpha_matting=True,
                        alpha_matting_background_threshold=value[0],
                        alpha_matting_foreground_threshold=value[1],
                        alpha_matting_erode_size=value[2],
                    )
            if isinstance(value, tuple):  # replace the background color
                if len(value) == 3:
                    value += (255,)
                img_update = remove(img_update, bgcolor=value)
            if isinstance(value, str):
                if confirm_rembg_models(value):
                    img_update = remove(img_update, session=new_session(value))
                else:
                    img_update = remove(img_update)
        elif "bgcolor" in k.lower():
            if isinstance(value, list):
                value = tuple(value)
            if isinstance(value, tuple):  # replace the background color
                if len(value) == 3:
                    value += (255,)
                img_update = remove(img_update, bgcolor=value)
    if filter_kws:
        for filter_name, filter_value in filter_kws.items():
            img_update = apply_filter(img_update, filter_name, filter_value)
    # Display the image if requested
    if show:
        if figsize is None:
            plt.figure(dpi=dpi)
        else:
            plt.figure(figsize=figsize, dpi=dpi)
        plt.imshow(img_update)
        plt.axis("on") if show_axis else plt.axis("off")
    return img_update


from sklearn.decomposition import PCA
from skimage import transform, feature, filters, measure
from skimage.color import rgb2gray
from scipy.fftpack import fftshift, fft2
import numpy as np
import cv2  # Used for template matching


def crop_black_borders(image):
    """Crop the black borders from a rotated image."""
    # Convert the image to grayscale if it's not already
    if image.ndim == 3:
        gray_image = color.rgb2gray(image)
    else:
        gray_image = image

    # Find all the non-black (non-zero) pixels
    mask = gray_image > 0  # Mask for non-black pixels (assuming black is zero)
    coords = np.column_stack(np.where(mask))

    # Get the bounding box of non-black pixels
    if coords.any():  # Check if there are any non-black pixels
        y_min, x_min = coords.min(axis=0)
        y_max, x_max = coords.max(axis=0)

        # Crop the image to the bounding box
        cropped_image = image[y_min : y_max + 1, x_min : x_max + 1]
    else:
        # If the image is completely black (which shouldn't happen), return the original image
        cropped_image = image

    return cropped_image


def detect_angle(image, by="median", template=None):
    """Detect the angle of rotation using various methods."""
    # Convert to grayscale
    gray_image = rgb2gray(image)

    # Detect edges using Canny edge detector
    edges = feature.canny(gray_image, sigma=2)

    # Use Hough transform to detect lines
    lines = transform.probabilistic_hough_line(edges)

    if not lines and any(["me" in by, "pca" in by]):
        print("No lines detected. Adjust the edge detection parameters.")
        return 0

    # Hough Transform-based angle detection (Median/Mean)
    if "me" in by:
        angles = []
        for line in lines:
            (x0, y0), (x1, y1) = line
            angle = np.arctan2(y1 - y0, x1 - x0) * 180 / np.pi
            if 80 < abs(angle) < 100:
                angles.append(angle)
        if not angles:
            return 0
        if "di" in by:
            median_angle = np.median(angles)
            rotation_angle = (
                90 - median_angle if median_angle > 0 else -90 - median_angle
            )

            return rotation_angle
        else:
            mean_angle = np.mean(angles)
            rotation_angle = 90 - mean_angle if mean_angle > 0 else -90 - mean_angle

            return rotation_angle

    # PCA-based angle detection
    elif "pca" in by:
        y, x = np.nonzero(edges)
        if len(x) == 0:
            return 0
        pca = PCA(n_components=2)
        pca.fit(np.vstack((x, y)).T)
        angle = np.arctan2(pca.components_[0, 1], pca.components_[0, 0]) * 180 / np.pi
        return angle

    # Gradient Orientation-based angle detection
    elif "gra" in by:
        gx, gy = np.gradient(gray_image)
        angles = np.arctan2(gy, gx) * 180 / np.pi
        hist, bin_edges = np.histogram(angles, bins=360, range=(-180, 180))
        return bin_edges[np.argmax(hist)]

    # Template Matching-based angle detection
    elif "temp" in by:
        if template is None:
            # Automatically extract a template from the center of the image
            height, width = gray_image.shape
            center_x, center_y = width // 2, height // 2
            size = (
                min(height, width) // 4
            )  # Size of the template as a fraction of image size
            template = gray_image[
                center_y - size : center_y + size, center_x - size : center_x + size
            ]
        best_angle = None
        best_corr = -1
        for angle in range(0, 180, 1):  # Checking every degree
            rotated_template = transform.rotate(template, angle)
            res = cv2.matchTemplate(gray_image, rotated_template, cv2.TM_CCOEFF)
            _, max_val, _, _ = cv2.minMaxLoc(res)
            if max_val > best_corr:
                best_corr = max_val
                best_angle = angle
        return best_angle

    # Image Moments-based angle detection
    elif "mo" in by:
        moments = measure.moments_central(gray_image)
        angle = (
            0.5
            * np.arctan2(2 * moments[1, 1], moments[0, 2] - moments[2, 0])
            * 180
            / np.pi
        )
        return angle

    # Fourier Transform-based angle detection
    elif "fft" in by:
        f = fft2(gray_image)
        fshift = fftshift(f)
        magnitude_spectrum = np.log(np.abs(fshift) + 1)
        rows, cols = magnitude_spectrum.shape
        r, c = np.unravel_index(np.argmax(magnitude_spectrum), (rows, cols))
        angle = np.arctan2(r - rows // 2, c - cols // 2) * 180 / np.pi
        return angle

    else:
        print(f"Unknown method {by}")
        return 0


def imgsets(img, **kwargs):
    """
    Apply various enhancements and filters to an image using PIL's ImageEnhance and ImageFilter modules.

    Args:
        img (PIL.Image): The input image.
        sets (dict): A dictionary specifying the enhancements, filters, and their parameters.
        show (bool): Whether to display the enhanced image.
        show_axis (bool): Whether to display axes on the image plot.
        size (tuple): The size of the thumbnail, cover, contain, or fit operation.
        dpi (int): Dots per inch for the displayed image.
        figsize (tuple): The size of the figure for displaying the image.
        auto (bool): Whether to automatically enhance the image based on its characteristics.

    Returns:
        PIL.Image: The enhanced image.

    Supported enhancements and filters:
        - "sharpness": Adjusts the sharpness of the image. Values > 1 increase sharpness, while values < 1 decrease sharpness.
        - "contrast": Adjusts the contrast of the image. Values > 1 increase contrast, while values < 1 decrease contrast.
        - "brightness": Adjusts the brightness of the image. Values > 1 increase brightness, while values < 1 decrease brightness.
        - "color": Adjusts the color saturation of the image. Values > 1 increase saturation, while values < 1 decrease saturation.
        - "rotate": Rotates the image by the specified angle.
        - "crop" or "cut": Crops the image. The value should be a tuple specifying the crop box as (left, upper, right, lower).
        - "size": Resizes the image to the specified dimensions.
        - "thumbnail": Resizes the image to fit within the given size while preserving aspect ratio.
        - "cover": Resizes and crops the image to fill the specified size.
        - "contain": Resizes the image to fit within the specified size, adding borders if necessary.
        - "fit": Resizes and pads the image to fit within the specified size.
        - "filter": Applies various filters to the image (e.g., BLUR, CONTOUR, EDGE_ENHANCE).

    Note:
        The "color" and "enhance" enhancements are not implemented in this function.
    """
    supported_filters = [
        "BLUR",
        "CONTOUR",
        "DETAIL",
        "EDGE_ENHANCE",
        "EDGE_ENHANCE_MORE",
        "EMBOSS",
        "FIND_EDGES",
        "SHARPEN",
        "SMOOTH",
        "SMOOTH_MORE",
        "MIN_FILTER",
        "MAX_FILTER",
        "MODE_FILTER",
        "MULTIBAND_FILTER",
        "GAUSSIAN_BLUR",
        "BOX_BLUR",
        "MEDIAN_FILTER",
    ]
    print('usage: imgsets(dir_img, contrast="auto", rm=True, color=2.2)')
    print("\nlog:\n")

    def confirm_rembg_models(model_name):
        models_support = [
            "u2net",
            "u2netp",
            "u2net_human_seg",
            "u2net_cloth_seg",
            "silueta",
            "isnet-general-use",
            "isnet-anime",
            "sam",
        ]
        if model_name in models_support:
            print(f"model_name: {model_name}")
            return model_name
        else:
            print(
                f"{model_name} cannot be found, check the name:{models_support}, default('isnet-general-use') has been used"
            )
            return "isnet-general-use"

    def auto_enhance(img):
        """
        Automatically enhances the image based on its characteristics.
        Args:
            img (PIL.Image): The input image.
        Returns:
            dict: A dictionary containing the optimal enhancement values.
        """
        # Determine the bit depth based on the image mode
        if img.mode in ["1", "L", "P", "RGB", "YCbCr", "LAB", "HSV"]:
            # 8-bit depth per channel
            bit_depth = 8
        elif img.mode in ["RGBA", "CMYK"]:
            # 8-bit depth per channel + alpha (RGBA) or additional channels (CMYK)
            bit_depth = 8
        elif img.mode in ["I", "F"]:
            # 16-bit depth per channel (integer or floating-point)
            bit_depth = 16
        else:
            raise ValueError("Unsupported image mode")
        # Calculate the brightness and contrast for each channel
        num_channels = len(img.getbands())
        brightness_factors = []
        contrast_factors = []
        for channel in range(num_channels):
            channel_histogram = img.split()[channel].histogram()
            brightness = sum(i * w for i, w in enumerate(channel_histogram)) / sum(
                channel_histogram
            )
            channel_min, channel_max = img.split()[channel].getextrema()
            contrast = channel_max - channel_min
            # Adjust calculations based on bit depth
            normalization_factor = 2**bit_depth - 1  # Max value for the given bit depth
            brightness_factor = (
                1.0 + (brightness - normalization_factor / 2) / normalization_factor
            )
            contrast_factor = (
                1.0 + (contrast - normalization_factor / 2) / normalization_factor
            )
            brightness_factors.append(brightness_factor)
            contrast_factors.append(contrast_factor)
        # Calculate the average brightness and contrast factors across channels
        avg_brightness_factor = sum(brightness_factors) / num_channels
        avg_contrast_factor = sum(contrast_factors) / num_channels
        return {"brightness": avg_brightness_factor, "contrast": avg_contrast_factor}

    # Load image if input is a file path
    if isinstance(img, str):
        img = load_img(img)
    img_update = img.copy()
    # Auto-enhance image if requested

    auto = kwargs.get("auto", False)
    show = kwargs.get("show", True)
    show_axis = kwargs.get("show_axis", False)
    size = kwargs.get("size", None)
    figsize = kwargs.get("figsize", None)
    dpi = kwargs.get("dpi", 100)

    if auto:
        kwargs = {**auto_enhance(img_update), **kwargs}

    for k, value in kwargs.items():
        if "shar" in k.lower():
            enhancer = ImageEnhance.Sharpness(img_update)
            img_update = enhancer.enhance(value)
        elif "col" in k.lower() and "bg" not in k.lower():
            enhancer = ImageEnhance.Color(img_update)
            img_update = enhancer.enhance(value)
        elif "contr" in k.lower():
            if value and isinstance(value, (float, int)):
                enhancer = ImageEnhance.Contrast(img_update)
                img_update = enhancer.enhance(value)
            else:
                print("autocontrasted")
                img_update = ImageOps.autocontrast(img_update)
        elif "bri" in k.lower():
            enhancer = ImageEnhance.Brightness(img_update)
            img_update = enhancer.enhance(value)
        elif "cro" in k.lower() or "cut" in k.lower():
            img_update = img_update.crop(value)
        elif "rota" in k.lower():
            if isinstance(value, str):
                value = detect_angle(img_update, by=value)
                print(f"rotated by {value}°")
            img_update = img_update.rotate(value)

        elif "si" in k.lower():
            img_update = img_update.resize(value)
        elif "thum" in k.lower():
            img_update.thumbnail(value)
        elif "cover" in k.lower():
            img_update = ImageOps.cover(img_update, size=value)
        elif "contain" in k.lower():
            img_update = ImageOps.contain(img_update, size=value)
        elif "fit" in k.lower():
            if isinstance(value, dict):
                for filter_name, filter_value in value.items():
                    img_update = apply_filter(img_update, filter_name, filter_value)
            else:
                img_update = ImageOps.fit(img_update, size=value)
        elif "pad" in k.lower():
            img_update = ImageOps.pad(img_update, size=value)
        elif "rem" in k.lower() or "rm" in k.lower() or "back" in k.lower():
            if isinstance(value, bool):
                session = new_session("isnet-general-use")
                img_update = remove(img_update, session=session)
            elif value and isinstance(value, (int, float, list)):
                print(
                    'example usage: {"rm":[alpha_matting_background_threshold(20),alpha_matting_foreground_threshold(270),alpha_matting_erode_sive(11)]}'
                )
                print("https://github.com/danielgatis/rembg/blob/main/USAGE.md")
                #     ###            Parameters:
                #         data (Union[bytes, PILImage, np.ndarray]): The input image data.
                #         alpha_matting (bool, optional): Flag indicating whether to use alpha matting. Defaults to False.
                #         alpha_matting_foreground_threshold (int, optional): Foreground threshold for alpha matting. Defaults to 240.
                #         alpha_matting_background_threshold (int, optional): Background threshold for alpha matting. Defaults to 10.
                #         alpha_matting_erode_size (int, optional): Erosion size for alpha matting. Defaults to 10.
                #         session (Optional[BaseSession], optional): A session object for the 'u2net' model. Defaults to None.
                #         only_mask (bool, optional): Flag indicating whether to return only the binary masks. Defaults to False.
                #         post_process_mask (bool, optional): Flag indicating whether to post-process the masks. Defaults to False.
                #         bgcolor (Optional[Tuple[int, int, int, int]], optional): Background color for the cutout image. Defaults to None.
                #  ###
                if isinstance(value, int):
                    value = [value]
                if len(value) < 2:
                    img_update = remove(
                        img_update,
                        alpha_matting=True,
                        alpha_matting_background_threshold=value,
                    )
                elif 2 <= len(value) < 3:
                    img_update = remove(
                        img_update,
                        alpha_matting=True,
                        alpha_matting_background_threshold=value[0],
                        alpha_matting_foreground_threshold=value[1],
                    )
                elif 3 <= len(value) < 4:
                    img_update = remove(
                        img_update,
                        alpha_matting=True,
                        alpha_matting_background_threshold=value[0],
                        alpha_matting_foreground_threshold=value[1],
                        alpha_matting_erode_size=value[2],
                    )
            elif isinstance(value, tuple):  # replace the background color
                if len(value) == 3:
                    value += (255,)
                img_update = remove(img_update, bgcolor=value)
            elif isinstance(value, str):
                if confirm_rembg_models(value):
                    img_update = remove(img_update, session=new_session(value))
                else:
                    img_update = remove(img_update)
        elif "bg" in k.lower() and "color" in k.lower():
            if isinstance(value, list):
                value = tuple(value)
            if isinstance(value, tuple):  # replace the background color
                if len(value) == 3:
                    value += (255,)
                img_update = remove(img_update, bgcolor=value)
    # Display the image if requested
    if show:
        if figsize is None:
            plt.figure(dpi=dpi)
        else:
            plt.figure(figsize=figsize, dpi=dpi)
        plt.imshow(img_update)
        if show_axis:
            plt.axis("on")  # Turn on axis
            plt.minorticks_on()
            plt.grid(
                which="both", linestyle="--", linewidth=0.5, color="gray", alpha=0.7
            )

        else:
            plt.axis("off")  # Turn off axis
    return img_update


# # usage:
# img = imgsets(
#     fpath,
#     sets={"rota": -5},
#     dpi=200,
#     filter_kws={"EMBOSS": 5, "sharpen": 5, "EDGE_ENHANCE_MORE": 10},
#     show_axis=True,
# )


def thumbnail(dir_img_list, figsize=(10, 10), dpi=100, dir_save=None, kind=".png"):
    """
    Display a thumbnail figure of all images in the specified directory.
    Args:
        dir_img_list (list): List of the Directory containing the images.
    """
    num_images = len(dir_img_list)
    if not kind.startswith("."):
        kind = "." + kind

    if num_images == 0:
        print("No images found to display.")
        return
    grid_size = int(num_images**0.5) + 1  # Determine grid size
    fig, axs = plt.subplots(grid_size, grid_size, figsize=figsize, dpi=dpi)
    for ax, image_file in zip(axs.flatten(), dir_img_list):
        try:
            img = Image.open(image_file)
            ax.imshow(img)
            ax.axis("off")
        except:
            continue
    # for ax in axs.flatten():
    #     ax.axis('off')
    [ax.axis("off") for ax in axs.flatten()]
    plt.tight_layout()
    if dir_save is None:
        plt.show()
    else:
        if basename(dir_save):
            fname = basename(dir_save) + kind
        else:
            fname = "_thumbnail_" + basename(dirname(dir_save)[:-1]) + ".png"
        if dirname(dir_img_list[0]) == dirname(dir_save):
            figsave(dirname(dir_save[:-1]), fname)
        else:
            figsave(dirname(dir_save), fname)


# usage:
# fpath = "/Users/macjianfeng/Dropbox/github/python/py2ls/tests/xample_netfinder/images/"
# thumbnail(listdir(fpath,'png').fpath.to_list(),dir_save=dirname(fpath))
def read_mplstyle(style_file):
    # Load the style file
    plt.style.use(style_file)

    # Get the current style properties
    style_dict = plt.rcParams

    # Convert to dictionary
    style_dict = dict(style_dict)
    # Print the style dictionary
    for i, j in style_dict.items():
        print(f"\n{i}::::{j}")
    return style_dict


# #example usage:
# style_file = "/ std-colors.mplstyle"
# style_dict = read_mplstyle(style_file)


# search and fine the director of the libary, which installed at local
def dir_lib(lib_oi):
    import site

    # Get the site-packages directory
    f = listdir(site.getsitepackages()[0], "folder")

    # Find Seaborn directory within site-packages
    dir_list = []
    for directory in f.fpath:
        if lib_oi in directory.lower():
            dir_list.append(directory)

    if dir_list != []:
        print(f"{lib_oi} directory:", dir_list)
    else:
        print(f"Cannot find the {lib_oi} in site-packages directory.")
    return dir_list


# example usage:
# dir_lib("seaborn")

""" 
    # n = 7
    # clist = get_color(n, cmap="auto", how="linspace")  # get_color(100)
    # plt.figure(figsize=[8, 5], dpi=100)
    # x = np.linspace(0, 2 * np.pi, 50) * 100
    # y = np.sin(x)
    # for i in range(1, n + 1):
    #     plt.plot(x, y + i, c=clist[i - 1], lw=5, label=str(i))
    # plt.legend()
    # plt.ylim(-2, 20)
    # figsets(plt.gca(), {"style": "whitegrid"}) """


class FileInfo:
    def __init__(
        self,
        size,
        creation_time,
        ctime,
        mod_time,
        mtime,
        parent_dir,
        fname,
        kind,
        extra_info=None,
    ):
        self.size = size
        self.creation_time = creation_time
        self.ctime = ctime
        self.mod_time = mod_time
        self.mtime = mtime
        self.parent_dir = parent_dir
        self.fname = fname
        self.kind = kind
        if extra_info:
            for key, value in extra_info.items():
                setattr(self, key, value)
        print("to show the res: 'finfo(fpath).show()'")

    def __repr__(self):
        return (
            f"FileInfo(size={self.size} MB, creation_time='{self.creation_time}', "
            f"ctime='{self.ctime}', mod_time='{self.mod_time}', mtime='{self.mtime}', "
            f"parent_dir='{self.parent_dir}', fname='{self.fname}', kind='{self.kind}')"
        )

    def __str__(self):
        return (
            f"FileInfo:\n"
            f"  Size: {self.size} MB\n"
            f"  Creation Time: {self.creation_time}\n"
            f"  CTime: {self.ctime}\n"
            f"  Modification Time: {self.mod_time}\n"
            f"  MTime: {self.mtime}\n"
            f"  Parent Directory: {self.parent_dir}\n"
            f"  File Name: {self.fname}\n"
            f"  Kind: {self.kind}"
        )

    def show(self):
        # Convert the object to a dictionary
        return {
            "size": self.size,
            "creation_time": self.creation_time,
            "ctime": self.ctime,
            "mod_time": self.mod_time,
            "mtime": self.mtime,
            "parent_dir": self.parent_dir,
            "fname": self.fname,
            "kind": self.kind,
            **{
                key: getattr(self, key)
                for key in vars(self)
                if key
                not in [
                    "size",
                    "creation_time",
                    "ctime",
                    "mod_time",
                    "mtime",
                    "parent_dir",
                    "fname",
                    "kind",
                ]
            },
        }


def finfo(fpath):
    fname, fmt = os.path.splitext(fpath)
    dir_par = os.path.dirname(fpath) + "/"
    data = {
        "size": round(os.path.getsize(fpath) / 1024 / 1024, 3),
        "creation_time": time.ctime(os.path.getctime(fpath)),
        "ctime": time.ctime(os.path.getctime(fpath)),
        "mod_time": time.ctime(os.path.getmtime(fpath)),
        "mtime": time.ctime(os.path.getmtime(fpath)),
        "parent_dir": dir_par,
        "fname": fname.replace(dir_par, ""),
        "kind": fmt,
    }
    extra_info = {}
    if data["kind"] == ".pdf":
        extra_info = pdfinfo_from_path(fpath)

    return FileInfo(
        size=data["size"],
        creation_time=data["creation_time"],
        ctime=data["ctime"],
        mod_time=data["mod_time"],
        mtime=data["mtime"],
        parent_dir=data["parent_dir"],
        fname=data["fname"],
        kind=data["kind"],
        extra_info=extra_info,
    )
